"""Generate the supervisor-meeting slide deck from live DB state.

Styled after `MSc Progress Slides 3.pptx`: 16:9 widescreen, dark hero
header bar on every slide, Calibri throughout, teal accent palette,
white cards with thin accent stripes. Every number on the deck is
queried from SQLite at generation time, so re-running this script
produces an up-to-date deck.

    uv run python scripts/generate_slides.py

Output: docs/PROGRESS_SLIDES_<YYYY-MM-DD>.pptx
"""

from __future__ import annotations

import argparse
import sqlite3
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
import sys
sys.path.insert(0, str(REPO_ROOT))
from dashboard.lib.currency import USD_TO_GBP

DB_PATH = REPO_ROOT / "data" / "odmi.db"
OUT_DIR = REPO_ROOT / "docs"


# ============================================================
# Palette and type system (mirrors MSc Progress Slides 3.pptx)
# ============================================================

DARK = RGBColor(0x1A, 0x20, 0x2C)
HEAD = RGBColor(0x1A, 0x20, 0x2C)
BODY = RGBColor(0x2D, 0x37, 0x48)
MUTED = RGBColor(0x71, 0x80, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xCB, 0xD5, 0xE0)
SURFACE = RGBColor(0xF7, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
TEAL_BRIGHT = RGBColor(0x14, 0xB8, 0xA6)
SUCCESS = RGBColor(0x38, 0xA1, 0x69)
WARNING = RGBColor(0xD6, 0x9E, 0x2E)
DANGER = RGBColor(0xC5, 0x30, 0x30)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)

FONT = "Calibri"


# ============================================================
# Data
# ============================================================

def fetch_stats() -> dict:
    with sqlite3.connect(DB_PATH) as conn:
        cur = conn.cursor()
        n_questions = cur.execute("SELECT COUNT(*) FROM questions").fetchone()[0]
        n_ground_truth = cur.execute(
            "SELECT COUNT(*) FROM ground_truth"
        ).fetchone()[0]
        n_researcher = cur.execute(
            "SELECT COUNT(*) FROM phase2_researcher_runs"
        ).fetchone()[0]
        n_verifier = cur.execute(
            "SELECT COUNT(*) FROM phase2_verifier_runs"
        ).fetchone()[0]
        n_adjudications = cur.execute(
            "SELECT COUNT(*) FROM phase2_adjudications"
        ).fetchone()[0]
        n_finals = cur.execute("SELECT COUNT(*) FROM phase2_final").fetchone()[0]

        # Per-country match vs differ vs ODMI ground truth.
        country_outcomes = cur.execute(
            """SELECT f.country_code,
                      SUM(CASE
                          WHEN LOWER(TRIM(f.final_answer)) =
                               LOWER(TRIM(gt.response))
                            OR (LOWER(TRIM(f.final_answer)) = 'yes'
                                AND LOWER(TRIM(gt.response)) LIKE 'yes%')
                          THEN 1 ELSE 0 END) AS matched,
                      SUM(CASE
                          WHEN gt.response IS NOT NULL
                               AND NOT (
                                 LOWER(TRIM(f.final_answer)) =
                                 LOWER(TRIM(gt.response))
                                 OR (LOWER(TRIM(f.final_answer)) = 'yes'
                                     AND LOWER(TRIM(gt.response)) LIKE 'yes%')
                               )
                          THEN 1 ELSE 0 END) AS differed
               FROM phase2_final f
               LEFT JOIN ground_truth gt
                  ON gt.question_id = f.question_id
                 AND gt.country_code = f.country_code
               GROUP BY f.country_code
               ORDER BY f.country_code"""
        ).fetchall()

        accuracy_row = cur.execute(
            """SELECT
                  SUM(CASE
                      WHEN LOWER(TRIM(f.final_answer)) =
                           LOWER(TRIM(gt.response))
                        OR (LOWER(TRIM(f.final_answer)) = 'yes'
                            AND LOWER(TRIM(gt.response)) LIKE 'yes%')
                      THEN 1 ELSE 0 END) AS matched,
                  COUNT(gt.response) AS comparable
               FROM phase2_final f
               LEFT JOIN ground_truth gt
                  ON gt.question_id = f.question_id
                 AND gt.country_code = f.country_code
               WHERE gt.response IS NOT NULL"""
        ).fetchone()
        n_match = int(accuracy_row[0] or 0)
        n_comparable = int(accuracy_row[1] or 0)
        accuracy = (n_match / n_comparable) if n_comparable > 0 else None

        cost_total = cur.execute(
            "SELECT COALESCE(SUM(estimated_cost_usd), 0) FROM claude_usage_log"
        ).fetchone()[0]

        avg_runtime = cur.execute(
            """SELECT AVG(
                 (julianday(ended_at) - julianday(started_at)) * 86400
               )
               FROM subtrio_status
               WHERE stage = 'done' AND ended_at IS NOT NULL"""
        ).fetchone()[0] or 0.0

    return {
        "n_questions": n_questions,
        "n_ground_truth": n_ground_truth,
        "n_researcher": n_researcher,
        "n_verifier": n_verifier,
        "n_adjudications": n_adjudications,
        "n_finals": n_finals,
        "country_outcomes": country_outcomes,
        "n_match": n_match,
        "n_comparable": n_comparable,
        "accuracy": accuracy,
        "cost_total": cost_total,
        "cost_total_gbp": (cost_total or 0.0) * USD_TO_GBP,
        "avg_runtime_s": avg_runtime,
    }


# ============================================================
# Low-level builders
# ============================================================

def set_text(
    tf, text: str, *,
    size: int = 12,
    bold: bool = False,
    colour: RGBColor = BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = FONT,
) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Emu(0)
    box.text_frame.margin_right = Emu(0)
    box.text_frame.margin_top = Emu(0)
    box.text_frame.margin_bottom = Emu(0)
    return box


def add_filled_rect(slide, x, y, w, h, *, fill: RGBColor, line_visible: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line_visible:
        shape.line.fill.background()
    return shape


def add_outlined_rect(
    slide, x, y, w, h, *,
    fill: RGBColor = WHITE,
    border: RGBColor = BORDER,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Emu(6350)  # ~0.5pt
    return shape


# ============================================================
# Slide template
# ============================================================

HEADER_HEIGHT = Inches(1.05)


def new_slide(prs: Presentation) -> object:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, 0, 0, prs.slide_width, HEADER_HEIGHT, fill=DARK)
    return slide


def header(slide, title: str, eyebrow: str | None = None) -> None:
    pad = Inches(0.5)
    if eyebrow:
        eyebrow_box = add_textbox(slide, pad, Inches(0.18), Inches(8.5), Inches(0.25))
        set_text(
            eyebrow_box.text_frame, eyebrow.upper(),
            size=9, bold=True, colour=TEAL_BRIGHT,
        )
        title_top = Inches(0.42)
    else:
        title_top = Inches(0.28)

    title_box = add_textbox(slide, pad, title_top, Inches(8.5), Inches(0.55))
    set_text(
        title_box.text_frame, title,
        size=24, bold=True, colour=WHITE,
    )


def page_footer(slide, prs: Presentation, page_num: int, total: int) -> None:
    pad = Inches(0.5)
    foot_y = prs.slide_height - Inches(0.32)
    fbox = add_textbox(slide, pad, foot_y, prs.slide_width - 2 * pad, Inches(0.2))
    set_text(
        fbox.text_frame,
        f"ODMI Agent Swarm  ·  Benjy Bream  ·  King's College London  ·  "
        f"{page_num} / {total}",
        size=8, colour=MUTED,
    )


# ============================================================
# Card builders
# ============================================================

def add_card(
    slide, x, y, w, h, *,
    title: str, body: str,
    accent: RGBColor = TEAL,
    title_size: int = 14,
    body_size: int = 11,
) -> None:
    add_outlined_rect(slide, x, y, w, h)
    add_filled_rect(slide, x, y, Inches(0.08), h, fill=accent)
    pad_x = Inches(0.25)
    pad_y = Inches(0.15)
    inner_x = x + Inches(0.18) + pad_x
    inner_w = w - Inches(0.18) - 2 * pad_x

    title_box = add_textbox(slide, inner_x, y + pad_y, inner_w, Inches(0.35))
    set_text(
        title_box.text_frame, title,
        size=title_size, bold=True, colour=HEAD,
    )

    body_box = add_textbox(
        slide, inner_x, y + pad_y + Inches(0.35),
        inner_w, h - pad_y - Inches(0.4),
    )
    set_text(body_box.text_frame, body, size=body_size, colour=BODY)


def add_kpi(
    slide, x, y, w, h, *,
    label: str, value: str, caption: str = "",
    accent: RGBColor = TEAL,
) -> None:
    add_outlined_rect(slide, x, y, w, h, fill=SURFACE)
    add_filled_rect(slide, x, y, w, Inches(0.06), fill=accent)

    pad_x = Inches(0.22)
    inner_w = w - 2 * pad_x

    # Label sits just under the accent stripe.
    label_box = add_textbox(slide, x + pad_x, y + Inches(0.18),
                            inner_w, Inches(0.22))
    set_text(label_box.text_frame, label.upper(),
             size=9, bold=True, colour=MUTED)

    # Value is the hero. Height matches the font's natural line so the
    # caption below doesn't get walked on.
    value_y = y + Inches(0.46)
    value_h = Inches(0.6)
    value_box = add_textbox(slide, x + pad_x, value_y, inner_w, value_h)
    set_text(value_box.text_frame, value,
             size=22, bold=True, colour=HEAD)

    if caption:
        cap_y = value_y + value_h + Inches(0.04)
        cap_h = max(Inches(0.2), y + h - cap_y - Inches(0.1))
        cap_box = add_textbox(slide, x + pad_x, cap_y, inner_w, cap_h)
        set_text(cap_box.text_frame, caption, size=9.5, colour=BODY)


def add_numbered_step(
    slide, x, y, w, h, *,
    number: int, title: str, body: str,
    tag: str | None = None,
    tag_colour: RGBColor = TEAL,
) -> None:
    add_outlined_rect(slide, x, y, w, h, fill=SURFACE)

    circle_d = Inches(0.5)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + Inches(0.2), y + (h - circle_d) / 2,
        circle_d, circle_d,
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK
    circle.line.fill.background()
    set_text(
        circle.text_frame, str(number),
        size=16, bold=True, colour=WHITE, align=PP_ALIGN.CENTER,
    )

    inner_x = x + Inches(0.85)
    inner_w = w - Inches(1.1)

    title_box = add_textbox(slide, inner_x, y + Inches(0.18), inner_w, Inches(0.3))
    set_text(title_box.text_frame, title, size=13, bold=True, colour=HEAD)

    body_box = add_textbox(slide, inner_x, y + Inches(0.5), inner_w, h - Inches(0.6))
    set_text(body_box.text_frame, body, size=10.5, colour=BODY)

    if tag:
        tag_w = Inches(0.7)
        tag_box = add_textbox(
            slide, x + w - tag_w - Inches(0.15),
            y + Inches(0.18), tag_w, Inches(0.25),
        )
        set_text(
            tag_box.text_frame, tag.upper(),
            size=9, bold=True, colour=tag_colour, align=PP_ALIGN.RIGHT,
        )


# ============================================================
# Slides
# ============================================================

def slide_title(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=DARK)

    add_filled_rect(slide, 0, Inches(2.0), Inches(0.6), Inches(2.0), fill=TEAL)

    eyebrow = add_textbox(slide, Inches(1.0), Inches(1.2), Inches(8), Inches(0.4))
    set_text(
        eyebrow.text_frame,
        "MSc ADVANCED COMPUTING · INDIVIDUAL PROJECT · STATUS BRIEFING",
        size=10, bold=True, colour=TEAL_BRIGHT,
    )

    title = add_textbox(slide, Inches(1.0), Inches(1.65), Inches(8.5), Inches(2.2))
    tf = title.text_frame
    tf.clear()
    for i, line in enumerate(["Agentic AI for Open", "Data Maturity", "Assessment"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

    sub = add_textbox(slide, Inches(1.0), Inches(4.0), Inches(8.5), Inches(0.6))
    set_text(
        sub.text_frame,
        "Status snapshot and next steps for supervisor review",
        size=14, colour=PALE,
    )

    today = datetime.now().strftime("%d %B %Y")
    foot = add_textbox(slide, Inches(1.0), Inches(5.0), Inches(8.5), Inches(0.4))
    set_text(
        foot.text_frame,
        f"Benjy Bream  ·  King's College London  ·  {today}",
        size=11, colour=WHITE,
    )


def slide_what_is_built(prs: Presentation, stats: dict) -> None:
    slide = new_slide(prs)
    header(slide, "What is built", eyebrow="Capability snapshot")

    y = Inches(1.25)
    h = Inches(1.5)
    col_w = Inches(2.225)
    gap = Inches(0.07)
    x0 = Inches(0.5)

    acc_str = (
        f"{stats['accuracy']:.0%}" if stats["accuracy"] is not None else "—"
    )
    add_kpi(slide, x0 + 0 * (col_w + gap), y, col_w, h,
            label="Pairs finalised",
            value=str(stats["n_finals"]),
            caption=f"{stats['n_researcher']} Researcher · {stats['n_verifier']} Verifier · "
                    f"{stats['n_adjudications']} Adjudicator.")
    add_kpi(slide, x0 + 1 * (col_w + gap), y, col_w, h,
            label="Ground-truth coverage",
            value=f"{stats['n_ground_truth']:,}",
            caption="ODMI answers loaded (36 countries × 143 questions).",
            accent=NAVY)
    add_kpi(slide, x0 + 2 * (col_w + gap), y, col_w, h,
            label="Accuracy vs ODMI",
            value=acc_str,
            caption=f"{stats['n_match']} / {stats['n_comparable']} match "
                    f"against ODMI 2025.",
            accent=SUCCESS)
    add_kpi(slide, x0 + 3 * (col_w + gap), y, col_w, h,
            label="Total LLM spend",
            value=f"£{stats['cost_total_gbp']:.2f}",
            caption="Via CLIProxyAPI on Claude Max.",
            accent=TEAL_BRIGHT)

    y2 = Inches(2.95)
    card_w = Inches(4.55)
    card_h = Inches(1.55)
    gap2 = Inches(0.15)
    add_card(slide, Inches(0.5), y2, card_w, card_h,
             title="Three-agent swarm, end-to-end",
             body=f"Researcher proposes, Verifier tries to disprove, "
                  f"Adjudicator steps in when retries are exhausted. "
                  f"Average runtime ~{stats['avg_runtime_s']:.0f}s.",
             accent=TEAL)
    add_card(slide, Inches(0.5) + card_w + gap2, y2, card_w, card_h,
             title="Live dashboard with ODMI ground truth",
             body="Streamlit app over SQLite. Every LLM call logs "
                  "tokens, latency, cost, prompt version. Each finalised "
                  "swarm pair is compared to ODMI's recorded answer "
                  "automatically (5,148 ground-truth rows loaded).",
             accent=NAVY)

    page_footer(slide, prs, 2, 8)


def slide_how_it_works(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "How the swarm answers one question",
           eyebrow="Coordinator state machine")

    row_y = Inches(1.6)
    row_h = Inches(1.65)
    box_w = Inches(2.7)
    gap = Inches(0.45)
    x_start = (prs.slide_width - (3 * box_w + 2 * gap)) // 2

    def agent_box(x, label, role, accent):
        add_outlined_rect(slide, x, row_y, box_w, row_h)
        add_filled_rect(slide, x, row_y, box_w, Inches(0.35), fill=accent)
        lab = add_textbox(slide, x + Inches(0.2), row_y + Inches(0.05),
                          box_w - Inches(0.4), Inches(0.27))
        set_text(lab.text_frame, label, size=12, bold=True, colour=WHITE)
        body = add_textbox(
            slide, x + Inches(0.2), row_y + Inches(0.5),
            box_w - Inches(0.4), row_h - Inches(0.6),
        )
        set_text(body.text_frame, role, size=10.5, colour=BODY)

    x1 = x_start
    x2 = x1 + box_w + gap
    x3 = x2 + box_w + gap

    agent_box(
        x1, "RESEARCHER",
        "Tavily search → fetch URL → reason → return "
        "{answer, source_url, evidence_quote, confidence}.",
        TEAL,
    )
    agent_box(
        x2, "VERIFIER",
        "Substring-checks the Researcher's quote, runs an independent "
        "search, and reasons adversarially (default strategy: find a "
        "way to disprove the claim). Returns pass / fail with "
        "counter-evidence.",
        NAVY,
    )
    agent_box(
        x3, "ADJUDICATOR",
        "Fires only on retry exhaustion. Reads the full history. "
        "Picks a winner or escalates to a human queue.",
        DANGER,
    )

    arrow_y = row_y + row_h // 2
    for left, right in [(x1 + box_w, x2), (x2 + box_w, x3)]:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            left + Inches(0.05), arrow_y - Inches(0.12),
            gap - Inches(0.1), Inches(0.24),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()

    callout_y = row_y + row_h + Inches(0.35)
    callout_w = (prs.slide_width - Inches(1.0) - Inches(0.3)) // 2
    callout_h = Inches(0.95)

    # Left callout: termination rule.
    add_outlined_rect(
        slide, Inches(0.5), callout_y, callout_w, callout_h,
        fill=SURFACE, border=TEAL,
    )
    cap_l = add_textbox(
        slide, Inches(0.7), callout_y + Inches(0.12),
        callout_w - Inches(0.4), Inches(0.32),
    )
    set_text(cap_l.text_frame, "TERMINATION RULE",
             size=9, bold=True, colour=TEAL)
    body_l = add_textbox(
        slide, Inches(0.7), callout_y + Inches(0.42),
        callout_w - Inches(0.4), Inches(0.5),
    )
    set_text(
        body_l.text_frame,
        "Researcher retries up to 3× on Verifier rejection. After "
        "that, Adjudicator decides. Confidence below 0.6 escalates "
        "to a human queue.",
        size=10, colour=BODY,
    )

    # Right callout: data-leakage guardrail (D24).
    right_x = Inches(0.5) + callout_w + Inches(0.3)
    add_outlined_rect(
        slide, right_x, callout_y, callout_w, callout_h,
        fill=SURFACE, border=DANGER,
    )
    cap_r = add_textbox(
        slide, right_x + Inches(0.2), callout_y + Inches(0.12),
        callout_w - Inches(0.4), Inches(0.32),
    )
    set_text(cap_r.text_frame, "DATA-LEAKAGE GUARDRAIL (D24)",
             size=9, bold=True, colour=DANGER)
    body_r = add_textbox(
        slide, right_x + Inches(0.2), callout_y + Inches(0.42),
        callout_w - Inches(0.4), Inches(0.5),
    )
    set_text(
        body_r.text_frame,
        "ODMI publications (data.europa.eu and mirrors) banned at "
        "five layers: search, fetch, validator, prompts, audit. "
        "An audit script catches anything that slips through.",
        size=10, colour=BODY,
    )

    page_footer(slide, prs, 3, 8)


def slide_dashboard_highlights(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Best bits of the dashboard",
           eyebrow="What I'll demo live")

    items = [
        ("Run Console",
         "Multiselect questions × countries, pick the Verifier "
         "strategy and the model triple, dispatch a batch. Cost is "
         "estimated before the run starts.",
         "▶", TEAL),
        ("Results · Cards view with ODMI side-by-side",
         "One card per finalised pair: question, swarm answer with "
         "evidence quote, clickable source URL, plus ODMI's recorded "
         "answer and explanation. Match / differ badge at a glance.",
         "📋", NAVY),
        ("Accuracy chart vs ODMI",
         "Per-country stacked bar of pairs that match ODMI 2025 versus "
         "those that differ. Headline accuracy KPI is read straight off "
         "the comparison.",
         "📊", SUCCESS),
        ("Cost tracking",
         "Costs page shows the 5-hour rolling spend, per-model "
         "breakdown, and per-strategy averages. Rolling 5-hour spend and "
         "rate-limit hits on every page's sidebar.",
         "💷", TEAL_DARK),
    ]

    y = Inches(1.35)
    w = Inches(4.55)
    h = Inches(1.85)
    gap = Inches(0.15)

    for i, (title, body, glyph, accent) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (w + gap)
        yy = y + row * (h + gap)

        add_outlined_rect(slide, x, yy, w, h)
        add_filled_rect(slide, x, yy, Inches(0.08), h, fill=accent)

        # Glyph bubble.
        bubble_d = Inches(0.55)
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.3), yy + Inches(0.25),
            bubble_d, bubble_d,
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = accent
        bubble.line.fill.background()
        set_text(
            bubble.text_frame, glyph,
            size=18, bold=True, colour=WHITE, align=PP_ALIGN.CENTER,
        )

        # Title.
        t_box = add_textbox(
            slide, x + Inches(1.0), yy + Inches(0.28),
            w - Inches(1.2), Inches(0.4),
        )
        set_text(t_box.text_frame, title,
                 size=15, bold=True, colour=HEAD)

        # Body.
        b_box = add_textbox(
            slide, x + Inches(1.0), yy + Inches(0.7),
            w - Inches(1.2), h - Inches(0.8),
        )
        set_text(b_box.text_frame, body, size=10.5, colour=BODY)

    page_footer(slide, prs, 6, 8)


def slide_data_leakage(prs: Presentation) -> None:
    """Five-layer defence against ODMI source leakage (SPEC D24)."""
    slide = new_slide(prs)
    header(
        slide, "Why ODMI publications are banned",
        eyebrow="Data-leakage controls · D24",
    )

    intro = add_textbox(
        slide, Inches(0.5), Inches(1.25),
        prs.slide_width - Inches(1.0), Inches(0.45),
    )
    set_text(
        intro.text_frame,
        "The swarm is validated against ODMI's published "
        "merged_responses. Any URL that hosts, mirrors, or summarises "
        "ODMI's own answers leaks the signal we are predicting, so it "
        "must never reach the agents. Five layers enforce this:",
        size=11, colour=BODY,
    )

    row_y = Inches(1.85)
    row_h = Inches(2.4)
    n_boxes = 5
    margin = Inches(0.5)
    available = prs.slide_width - 2 * margin
    gap = Inches(0.18)
    box_w = (available - (n_boxes - 1) * gap) // n_boxes

    layers = [
        ("1", "SEARCH", TEAL,
         "Tavily exclude_domains · Brave −site: · post-filter drops "
         "any blocked URL before it reaches an agent."),
        ("2", "FETCH", TEAL_DARK,
         "fetch_text and head_ok refuse blocked URLs with "
         "failure_mode=\"blocked_data_leakage\". No network call."),
        ("3", "VALIDATOR", NAVY,
         "trust_score forced to 0.0 for blocked URLs. "
         "data.europa.eu removed from the trusted domain seed list."),
        ("4", "PROMPTS", WARNING,
         "Researcher V2, all four Verifier V2 strategies, and "
         "Adjudicator V2 carry an explicit forbidden_odmi_source rule."),
        ("5", "AUDIT", DANGER,
         "scripts/check_data_leakage.py scans every URL column in "
         "the swarm tables. Exits non-zero on any hit. --purge "
         "deletes tainted pair_runs."),
    ]

    for i, (number, name, accent, body_text) in enumerate(layers):
        x = margin + i * (box_w + gap)
        add_outlined_rect(slide, x, row_y, box_w, row_h)
        # Numbered top stripe.
        add_filled_rect(slide, x, row_y, box_w, Inches(0.35),
                        fill=accent)
        num_box = add_textbox(
            slide, x + Inches(0.2), row_y + Inches(0.05),
            box_w - Inches(0.4), Inches(0.27),
        )
        set_text(num_box.text_frame, f"{number}.  {name}",
                 size=11, bold=True, colour=WHITE)
        body = add_textbox(
            slide, x + Inches(0.15), row_y + Inches(0.5),
            box_w - Inches(0.3), row_h - Inches(0.6),
        )
        set_text(body.text_frame, body_text, size=9.5, colour=BODY)

    # Bottom callout — why five layers.
    callout_y = row_y + row_h + Inches(0.2)
    add_outlined_rect(
        slide, Inches(0.5), callout_y,
        prs.slide_width - Inches(1.0), Inches(0.55),
        fill=SURFACE, border=DANGER,
    )
    cap = add_textbox(
        slide, Inches(0.7), callout_y + Inches(0.1),
        prs.slide_width - Inches(1.4), Inches(0.4),
    )
    set_text(
        cap.text_frame,
        "Each layer fails in a different way. The audit catches whatever "
        "the four runtime layers miss; 30 historical violations were "
        "purged when D24 landed.",
        size=10.5, colour=BODY,
    )

    page_footer(slide, prs, 4, 8)


def slide_country_chart(prs: Presentation, stats: dict) -> None:
    slide = new_slide(prs)
    header(slide, "Results so far, per country",
           eyebrow="Finalised pairs · success vs failure")

    outcomes = stats["country_outcomes"]

    chart_box_x = Inches(0.5)
    chart_box_y = Inches(1.3)
    chart_box_w = prs.slide_width - Inches(1.0)
    chart_box_h = Inches(2.7)

    add_outlined_rect(slide, chart_box_x, chart_box_y, chart_box_w, chart_box_h,
                      fill=WHITE, border=BORDER)

    if not outcomes:
        msg = add_textbox(slide, chart_box_x + Inches(0.3),
                          chart_box_y + Inches(1.2),
                          chart_box_w - Inches(0.6), Inches(0.5))
        set_text(
            msg.text_frame,
            "No finalised pairs yet. Once the Coordinator writes "
            "phase2_final rows, this chart populates.",
            size=14, colour=BODY,
        )
        page_footer(slide, prs, 5, 8)
        return

    max_total = max(s + f for _, s, f in outcomes) or 1

    plot_left = chart_box_x + Inches(0.8)
    plot_top = chart_box_y + Inches(0.3)
    plot_width = chart_box_w - Inches(1.3)
    plot_height = chart_box_h - Inches(0.95)

    bar_slot = plot_width / len(outcomes)
    bar_w = Emu(int(bar_slot * 0.55))
    gap = (bar_slot - bar_w) // 2

    axis = slide.shapes.add_connector(
        1, plot_left - Inches(0.05), plot_top,
        plot_left - Inches(0.05), plot_top + plot_height,
    )
    axis.line.color.rgb = BORDER
    axis.line.width = Emu(6350)

    base = slide.shapes.add_connector(
        1, plot_left - Inches(0.05), plot_top + plot_height,
        plot_left + plot_width, plot_top + plot_height,
    )
    base.line.color.rgb = BORDER
    base.line.width = Emu(6350)

    for tick in range(0, max_total + 1):
        ratio = tick / max_total
        ty = plot_top + plot_height - Emu(int(plot_height * ratio))
        lbl = add_textbox(
            slide, plot_left - Inches(0.55), ty - Inches(0.1),
            Inches(0.45), Inches(0.2),
        )
        set_text(
            lbl.text_frame, str(tick),
            size=9, colour=MUTED, align=PP_ALIGN.RIGHT,
        )

    for i, (country, matched, differed) in enumerate(outcomes):
        total = matched + differed
        if total == 0:
            continue
        bar_x = plot_left + Emu(int(bar_slot * i + gap))

        if differed > 0:
            d_h = Emu(int(plot_height * differed / max_total))
            d_top = plot_top + plot_height - d_h
            add_filled_rect(slide, bar_x, d_top, bar_w, d_h, fill=DANGER)

        if matched > 0:
            d_h = Emu(int(plot_height * differed / max_total))
            m_h = Emu(int(plot_height * matched / max_total))
            m_top = plot_top + plot_height - d_h - m_h
            add_filled_rect(slide, bar_x, m_top, bar_w, m_h, fill=SUCCESS)

        lab_y = plot_top + plot_height + Inches(0.08)
        lab = add_textbox(
            slide, bar_x - Inches(0.4), lab_y,
            bar_w + Inches(0.8), Inches(0.25),
        )
        set_text(
            lab.text_frame, country,
            size=12, bold=True, colour=HEAD, align=PP_ALIGN.CENTER,
        )
        val = add_textbox(
            slide, bar_x - Inches(0.4), lab_y + Inches(0.22),
            bar_w + Inches(0.8), Inches(0.22),
        )
        set_text(
            val.text_frame, f"{matched} / {total}",
            size=9, colour=MUTED, align=PP_ALIGN.CENTER,
        )

    # Legend sits below the chart card, not inside it.
    leg_y = chart_box_y + chart_box_h + Inches(0.18)
    leg_x = chart_box_x + Inches(0.1)
    add_filled_rect(slide, leg_x, leg_y, Inches(0.18), Inches(0.18), fill=SUCCESS)
    leg1 = add_textbox(slide, leg_x + Inches(0.28), leg_y - Inches(0.03),
                       Inches(3.6), Inches(0.25))
    set_text(leg1.text_frame, "Matches ODMI 2025",
             size=10, colour=BODY)
    add_filled_rect(slide, leg_x + Inches(3.6), leg_y, Inches(0.18), Inches(0.18),
                    fill=DANGER)
    leg2 = add_textbox(slide, leg_x + Inches(3.88), leg_y - Inches(0.03),
                       Inches(4.0), Inches(0.25))
    set_text(leg2.text_frame, "Differs from ODMI 2025",
             size=10, colour=BODY)

    note_y = leg_y + Inches(0.45)
    add_filled_rect(slide, chart_box_x, note_y, Inches(0.08), Inches(0.6),
                    fill=WARNING)
    nb = add_textbox(
        slide, chart_box_x + Inches(0.2), note_y + Inches(0.04),
        chart_box_w - Inches(0.3), Inches(0.55),
    )
    set_text(
        nb.text_frame,
        "Comparison is to ODMI's 2025 merged_responses. The current "
        "sample is small and Policy-dimension heavy. ODMI's assessments "
        "are one cycle old, so where the swarm differs that does not "
        "automatically mean the swarm is wrong: the country may have "
        "changed since the assessment. Each disagreement is worth a "
        "human glance.",
        size=10, colour=BODY,
    )

    page_footer(slide, prs, 5, 8)


def slide_swarm_in_action(prs: Presentation) -> None:
    _img = REPO_ROOT / "docs" / "assets" / "swarm_in_action.png"
    if not _img.exists():
        return
    slide = new_slide(prs)
    header(slide, "Run Console — live view while the swarm runs",
           eyebrow="Swarm in action")
    slide.shapes.add_picture(
        str(_img),
        Inches(0.5), Inches(1.15),
        prs.slide_width - Inches(1.0), Inches(4.15),
    )
    page_footer(slide, prs, 7, 8)


def slide_next_steps(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Where this goes next",
           eyebrow="Short term and long term")

    col_w = Inches(4.55)
    gap = Inches(0.2)
    x_l = Inches(0.5)
    x_r = x_l + col_w + gap
    y = Inches(1.35)
    h_col = Inches(3.85)

    def column(x, header_label, header_fill, items):
        add_outlined_rect(slide, x, y, col_w, h_col, fill=SURFACE)
        add_filled_rect(slide, x, y, col_w, Inches(0.45), fill=header_fill)
        head_box = add_textbox(slide, x + Inches(0.25), y + Inches(0.1),
                               col_w - Inches(0.5), Inches(0.3))
        set_text(head_box.text_frame, header_label,
                 size=11, bold=True, colour=WHITE)

        item_h = (h_col - Inches(0.45) - Inches(0.2)) / len(items)
        row_y = y + Inches(0.55)
        for title, body in items:
            num_x = x + Inches(0.2)
            num_w = Inches(0.4)
            t_box = add_textbox(slide, num_x + num_w + Inches(0.05), row_y,
                                col_w - num_w - Inches(0.3), Inches(0.28))
            set_text(t_box.text_frame, title, size=11, bold=True, colour=HEAD)

            b_box = add_textbox(slide, num_x + num_w + Inches(0.05),
                                row_y + Inches(0.28),
                                col_w - num_w - Inches(0.3),
                                item_h - Inches(0.32))
            set_text(b_box.text_frame, body, size=10, colour=BODY)

            # Accent bullet on the left.
            add_filled_rect(slide, num_x, row_y + Inches(0.08),
                            Inches(0.18), Inches(0.18), fill=TEAL)
            row_y += item_h

    short_items = [
        ("Scale to harder questions",
         "Run the swarm beyond Policy on the easy four countries: "
         "Portal, Quality, and Impact dimensions. Track where accuracy "
         "starts to fall."),
        ("Add the lower-resource countries",
         "Bring Hungary and Estonia into the regular sweep. Country-"
         "language fetch is the biggest unknown."),
        ("Verifier strategy comparison",
         "Four strategies on the same Researcher rows. Hallucination "
         "catch rate, false rejection rate, tokens per Verifier run."),
    ]

    long_items = [
        ("Scale to all six countries",
         "Run all 143 questions across FR, DE, NL, RO, HU and EE — "
         "the full Phase B matrix."),
        ("Model cost-effectiveness experiments",
         "Compare Haiku, Sonnet and Opus on accuracy and cost across "
         "the matrix. Build the cost-accuracy surface and identify the "
         "optimal strategy."),
        ("July — dissertation write-up",
         "Analyse results, write methodology and findings chapters. "
         "Submission deadline: 2 August 2026."),
    ]

    column(x_l, "NEXT TWO WEEKS", DARK, short_items)
    column(x_r, "JUNE AND JULY", NAVY, long_items)

    page_footer(slide, prs, 8, 8)


# ============================================================
# PhD-seminar talk (5 slides, opinion-seeking)
# ============================================================
#
# A separate, tight deck for the Monday PhD-student seminar. The
# supervisor introduces the problem, so this opens on the system and
# spends its weight on method and on the open questions we want the
# room to argue. Numbers that are not in the live DB (EXP-1, the Malta
# baseline, the catalogue recompute) are committed, documented results,
# pinned here as constants with their source in docs/EXPERIMENTS.md.

def _list_column(
    slide, x, y, col_w, h_col, *,
    header_label: str, header_fill: RGBColor,
    items, bullet: RGBColor = TEAL,
) -> None:
    add_outlined_rect(slide, x, y, col_w, h_col, fill=SURFACE)
    add_filled_rect(slide, x, y, col_w, Inches(0.45), fill=header_fill)
    head_box = add_textbox(slide, x + Inches(0.25), y + Inches(0.1),
                           col_w - Inches(0.5), Inches(0.3))
    set_text(head_box.text_frame, header_label, size=11, bold=True, colour=WHITE)

    item_h = (h_col - Inches(0.45) - Inches(0.2)) / len(items)
    row_y = y + Inches(0.55)
    for title, body in items:
        num_x = x + Inches(0.2)
        num_w = Inches(0.4)
        t_box = add_textbox(slide, num_x + num_w + Inches(0.05), row_y,
                            col_w - num_w - Inches(0.3), Inches(0.28))
        set_text(t_box.text_frame, title, size=11, bold=True, colour=HEAD)
        b_box = add_textbox(slide, num_x + num_w + Inches(0.05),
                            row_y + Inches(0.28),
                            col_w - num_w - Inches(0.3), item_h - Inches(0.32))
        set_text(b_box.text_frame, body, size=10, colour=BODY)
        add_filled_rect(slide, num_x, row_y + Inches(0.08),
                        Inches(0.18), Inches(0.18), fill=bullet)
        row_y += item_h


def phd_slide_title(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=DARK)
    add_filled_rect(slide, 0, Inches(2.0), Inches(0.6), Inches(2.0), fill=TEAL)

    eyebrow = add_textbox(slide, Inches(1.0), Inches(1.2), Inches(8.5), Inches(0.4))
    set_text(
        eyebrow.text_frame,
        "MSc ADVANCED COMPUTING · PhD SEMINAR · FOR DISCUSSION",
        size=10, bold=True, colour=TEAL_BRIGHT,
    )

    title = add_textbox(slide, Inches(1.0), Inches(1.65), Inches(8.5), Inches(2.2))
    tf = title.text_frame
    tf.clear()
    for i, line in enumerate(["Agentic AI for Open", "Data Maturity", "Assessment"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

    sub = add_textbox(slide, Inches(1.0), Inches(4.0), Inches(8.5), Inches(0.6))
    set_text(
        sub.text_frame,
        "Method, experiments, and the open questions I want your view on",
        size=14, colour=PALE,
    )

    today = datetime.now().strftime("%d %B %Y")
    foot = add_textbox(slide, Inches(1.0), Inches(5.0), Inches(8.5), Inches(0.4))
    set_text(
        foot.text_frame,
        f"Benjy Bream  ·  King's College London  ·  {today}",
        size=11, colour=WHITE,
    )


def phd_slide_system(prs: Presentation, stats: dict) -> None:
    slide = new_slide(prs)
    header(slide, "An adversarial swarm, scored against ODMI's own answers",
           eyebrow="System and honest status")

    row_y = Inches(1.3)
    row_h = Inches(1.4)
    box_w = Inches(2.7)
    gap = Inches(0.45)
    x_start = (prs.slide_width - (3 * box_w + 2 * gap)) // 2

    def agent_box(x, label, role, accent):
        add_outlined_rect(slide, x, row_y, box_w, row_h)
        add_filled_rect(slide, x, row_y, box_w, Inches(0.35), fill=accent)
        lab = add_textbox(slide, x + Inches(0.2), row_y + Inches(0.05),
                          box_w - Inches(0.4), Inches(0.27))
        set_text(lab.text_frame, label, size=12, bold=True, colour=WHITE)
        body = add_textbox(slide, x + Inches(0.2), row_y + Inches(0.5),
                           box_w - Inches(0.4), row_h - Inches(0.6))
        set_text(body.text_frame, role, size=10, colour=BODY)

    x1 = x_start
    x2 = x1 + box_w + gap
    x3 = x2 + box_w + gap
    agent_box(x1, "RESEARCHER",
              "Searches, fetches, reasons. Returns answer + source URL "
              "+ quoted passage + confidence.", TEAL)
    agent_box(x2, "VERIFIER",
              "Independent search, prompted to disprove. Returns pass / "
              "fail with counter-evidence.", NAVY)
    agent_box(x3, "ADJUDICATOR",
              "Fires on retry exhaustion. Reads the full history, commits "
              "above the 0.65 floor or abstains.", DANGER)

    arrow_y = row_y + row_h // 2
    for left, right in [(x1 + box_w, x2), (x2 + box_w, x3)]:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            left + Inches(0.05), arrow_y - Inches(0.12),
            gap - Inches(0.1), Inches(0.24),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()

    add_card(
        slide, Inches(0.5), Inches(3.05), prs.slide_width - Inches(1.0),
        Inches(1.65),
        title="What is actually run, and what is not",
        body="Plain Python state machine, Claude via CLIProxyAPI. Every "
             "call logs model + prompt version + raw response + cost, so "
             f"an examiner replays it from logs alone. Run and adjudicated: "
             f"EXP-1 (DIY vs Tavily), the Malta base-rate baseline (60/60), "
             f"and the independent catalogue recompute. Everything else is "
             f"pre-registered and built, gated on search and Claude quota. "
             f"{stats['n_finals']} pairs finalised, "
             f"{stats['n_ground_truth']:,} ODMI ground-truth rows loaded.",
        accent=TEAL,
    )

    page_footer(slide, prs, 2, 6)


def phd_slide_impartiality(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Keeping the evaluation honest",
           eyebrow="Impartiality designed in, not asserted")

    col_w = Inches(4.45)
    gap = Inches(0.1)
    x_l = Inches(0.5)
    x_r = x_l + col_w + gap
    y1 = Inches(1.25)
    y2 = Inches(3.0)
    h = Inches(1.6)

    add_card(slide, x_l, y1, col_w, h,
             title="Pre-registered before the data (R1)",
             body="Each design is fixed in a dated git commit whose "
                  "timestamp predates its result file. The run decides "
                  "the answer, nothing else. No reverse-fitting.",
             accent=TEAL)
    add_card(slide, x_r, y1, col_w, h,
             title="Blind, swapped, cross-checked judge (R5/R6)",
             body="Evidence blinded and reduced to the bare domain, each "
                  "pair judged twice with arms swapped at temperature 0. A "
                  "second model family re-rates a seeded subsample.",
             accent=NAVY)
    add_card(slide, x_l, y2, col_w, h,
             title="Deny-list, equal and pre-fetch (R10 / D24)",
             body="ODMI publishes its own answers on data.europa.eu. That "
                  "domain is blocked for every arm before retrieval, so no "
                  "arm competes for a different set of usable sources.",
             accent=TEAL_DARK)
    add_card(slide, x_r, y2, col_w, h,
             title="The base-rate rule (R4) — the France lesson",
             body="France binary gold is 119 yes to 1 no, so 99% accuracy "
                  "measures nothing. Balance-aware metrics, and Malta "
                  "promoted to primary; France barred as a primary set.",
             accent=DANGER)

    page_footer(slide, prs, 3, 6)


def phd_slide_experiments(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Experiments: two real results, the rest queued",
           eyebrow="Findings and apparatus")

    col_w = Inches(4.45)
    gap = Inches(0.1)
    x_l = Inches(0.5)
    x_r = x_l + col_w + gap
    y1 = Inches(1.25)
    h1 = Inches(1.95)

    add_card(slide, x_l, y1, col_w, h1,
             title="EXP-1 · DIY search beats Tavily (run)",
             body="90 French pairs, blind position-swapped Opus judge. DIY "
                  "wins 89% of the 55 decided pairs (Wilson [78, 95], "
                  "p < 1e-4); 93% under strict exclusion. Cross-family "
                  "Mistral re-judge: Krippendorff alpha 0.648, which rebuts "
                  "the same-family self-preference worry.",
             accent=SUCCESS)
    add_card(slide, x_r, y1, col_w, h1,
             title="Catalogue recompute · self-report ceiling (run)",
             body="A deterministic DCAT-AP / SHACL tool recomputes the "
                  "Quality questions independently of the deny-listed MQA. "
                  "France self-reported >90% licence coverage; the "
                  "independent recompute reads ~38%. The self-report "
                  "ceiling, made measurable.",
             accent=NAVY)

    add_card(slide, Inches(0.5), Inches(3.45), prs.slide_width - Inches(1.0),
             Inches(1.3),
             title="Pre-registered, built, gated on quota",
             body="Malta baseline finalised 60/60. EXP-6 verifier-strategy "
                  "discrimination · EXP-7 retry chaining (flag-gated, "
                  "default off) · EXP-8/9 cost and model-variant surface · "
                  "EXP-10 Malta failure audit with a free confidence-floor "
                  "sweep. Code and pre-registration done; runs pending "
                  "search and Claude headroom.",
             accent=MUTED)

    page_footer(slide, prs, 4, 6)


def phd_slide_failure_modes(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Failure modes — fixed and still open",
           eyebrow="What broke, what remains")

    col_w = Inches(4.45)
    gap = Inches(0.1)
    x_l = Inches(0.5)
    x_r = x_l + col_w + gap
    y = Inches(1.25)
    h_col = Inches(3.55)

    fixed = [
        ("WAF 403 on data.gov.mt",
         "Cloudflare blocked head checks and killed valid answers. A "
         "Playwright render now retries on 403 / 429 / 503."),
        ("Resume stranded pairs",
         "The resume path reused failed Researcher rows, leaving 11 pairs "
         "stuck at 'researching'. Fixed in the coordinator."),
        ("Silent auth failure",
         "An empty auth token made every call fail as a misleading "
         "connection error. Fixed in the LLM client."),
    ]
    open_ = [
        ("Retrieval ceiling",
         "Self-report and deny-listed Quality questions are not "
         "answerable from the open web. The residual bottleneck."),
        ("Recall asymmetry",
         "Yes-gold recall 0.60 against no-gold 0.87: the swarm confirms a "
         "'no' more readily than a 'yes' on sparse evidence."),
        ("Thin-web countries",
         "Low-resource languages and sparse results untested at scale; "
         "language and difficulty are still confounded."),
    ]

    _list_column(slide, x_l, y, col_w, h_col,
                 header_label="FIXED", header_fill=TEAL_DARK,
                 items=fixed, bullet=SUCCESS)
    _list_column(slide, x_r, y, col_w, h_col,
                 header_label="STILL OPEN", header_fill=WARNING,
                 items=open_, bullet=WARNING)

    page_footer(slide, prs, 5, 6)


def phd_slide_open_questions(prs: Presentation) -> None:
    slide = new_slide(prs)
    header(slide, "Where I want your opinion",
           eyebrow="Next steps and open questions")

    col_w = Inches(4.45)
    gap = Inches(0.1)
    x_l = Inches(0.5)
    x_r = x_l + col_w + gap
    y = Inches(1.25)
    h_col = Inches(3.55)

    next_steps = [
        ("Freeze a held-out test stratum",
         "Tune on a subset, report final accuracy on unseen countries "
         "against the contemporaneous web."),
        ("Run the pre-registered set",
         "EXP-6 / 7 / 8 / 9 / 10 once quota resets; build the "
         "accuracy-cost surface."),
        ("Six-country sweep",
         "Wealth x maturity matrix across all 143 questions, base-rate "
         "checked per cell."),
    ]
    questions = [
        ("What is the test set?",
         "Validation is 2025 ODMI. Held-out countries, 2024 (stale 2026 "
         "web), or matched leave-one-out pairs?"),
        ("How to control for language?",
         "Malta + Netherlands, France with injected inversions, or "
         "matched country pairs differing only in language?"),
        ("Is the 6-country matrix fair?",
         "Or does wealth x maturity leave language and base rate "
         "confounded?"),
        ("Blindspots and missing experiments?",
         "Verifier ablation, confidence calibration, a human-anchored "
         "judge — what else am I not seeing?"),
    ]

    _list_column(slide, x_l, y, col_w, h_col,
                 header_label="NEXT STEPS", header_fill=DARK,
                 items=next_steps, bullet=TEAL)
    _list_column(slide, x_r, y, col_w, h_col,
                 header_label="QUESTIONS FOR THE ROOM", header_fill=NAVY,
                 items=questions, bullet=TEAL_BRIGHT)

    page_footer(slide, prs, 6, 6)


def build_phd_deck(stats: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    phd_slide_title(prs)
    phd_slide_system(prs, stats)
    phd_slide_impartiality(prs)
    phd_slide_experiments(prs)
    phd_slide_failure_modes(prs)
    phd_slide_open_questions(prs)
    return prs


# ============================================================
# Main
# ============================================================

def build_deck(stats: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_title(prs)
    slide_what_is_built(prs, stats)
    slide_how_it_works(prs)
    slide_data_leakage(prs)
    slide_country_chart(prs, stats)
    slide_dashboard_highlights(prs)
    slide_swarm_in_action(prs)
    slide_next_steps(prs)
    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--phd", action="store_true",
        help="Build the 5-slide PhD-seminar talk (opinion-seeking) instead "
             "of the 8-slide supervisor briefing.",
    )
    args = parser.parse_args()

    stats = fetch_stats()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y-%m-%d")
    if args.phd:
        prs = build_phd_deck(stats)
        out_path = OUT_DIR / f"PHD_TALK_{stamp}.pptx"
    else:
        prs = build_deck(stats)
        out_path = OUT_DIR / f"PROGRESS_SLIDES_{stamp}.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")
    print(f"  Slides: {len(prs.slides)}")
    if not args.phd:
        print(f"  Finals on chart: {len(stats['country_outcomes'])} country/countries")


if __name__ == "__main__":
    main()
